import streamlit as st
from google import genai
from google.genai import types
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib import styles
from PIL import Image
import io
import tempfile
import os

# ================================== #
# KONFIGURASI #
# ================================== #
st.set_page_config(page_title="AI Presentasi Modul Ajar Katolik", layout="wide")
st.title("⛪ AI Presentasi Modul Ajar Katolik")
st.write("Upload modul ajar DOCX/PDF → AI membuat presentasi interaktif otomatis")

api_key = st.sidebar.text_input("Gemini API Key", type="password")
if not api_key:
    st.warning("Masukkan Gemini API key")
    st.stop()

# Inisialisasi client Gemini
client = genai.Client(api_key=api_key)

# ================================== #
# WARNA TEMA KATOLIK #
# ================================== #
THEME = {"primary": "1F4E79", "secondary": "D4AF37", "background": "FFFFFF"}

# ================================== #
# BACA FILE #
# ================================== #
def read_docx(file):
    doc = Document(file)
    text = []
    for p in doc.paragraphs:
        if p.text.strip():
            text.append(p.text)
    return "\n".join(text)

def read_pdf(file):
    reader = PdfReader(file)
    text = []
    for page in reader.pages:
        try:
            text.append(page.extract_text())
        except:
            pass
    return "\n".join(text)

# ================================== #
# AI BUAT STRUKTUR SLIDE #
# ================================== #
def generate_slides(text):
    prompt = f"""
    Dari modul berikut: {text}
    Buat presentasi:
    Slide 1 Judul
    Slide 2 Tujuan Pembelajaran
    Slide 3 Materi
    Slide 4 Aktivitas
    Slide 5 Refleksi
    Slide 6 Asesmen
    Slide 7 Kesimpulan
    Maksimal 5 poin per slide
    """
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    return response.text

# ================================== #
# AI BUAT KUIS #
# ================================== #
def generate_quiz(text):
    prompt = f"""
    Buat 10 soal pilihan ganda dari: {text}
    format:
    Soal:
    A.
    B.
    C.
    D.
    Jawaban:
    """
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    return response.text

# ================================== #
# BUAT GAMBAR AI (Gemini Imagen) #
# ================================== #
def generate_image(topic):
    response = client.models.generate_images(
        model="imagen-3.0-generate-002",
        prompt=f"Ilustrasi pendidikan katolik: {topic}",
        config=types.GenerateImagesConfig(number_of_images=1, aspect_ratio="1:1")
    )
    return response.generated_images[0].image.image_bytes

# ================================== #
# SIMPAN GAMBAR KE TEMP #
# ================================== #
def save_image(image_bytes):
    path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
    with open(path, "wb") as f:
        f.write(image_bytes)
    return path

# ================================== #
# BUAT PPT #
# ================================== #
def create_ppt(content, quiz):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Presentasi Modul Ajar"
    title_slide.placeholders[1].text = "Pendidikan Agama Katolik"
    
    slides = content.split("Slide")
    for item in slides:
        if len(item) < 10:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        parts = item.split("\n")
        title = parts[0]
        slide.shapes.title.text = title
        
        textbox = slide.shapes.add_textbox(Inches(.5), Inches(1.2), Inches(5), Inches(4))
        frame = textbox.text_frame
        for line in parts[1:]:
            if line.strip():
                p = frame.add_paragraph()
                p.text = "• " + line
        try:
            image_bytes = generate_image(title)
            img = save_image(image_bytes)
            slide.shapes.add_picture(img, Inches(7), Inches(1.5), width=Inches(3))
        except:
            pass
            
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Kuis"
    slide.placeholders[1].text = quiz
    
    ppt_file = "hasil_presentasi.pptx"
    prs.save(ppt_file)
    return ppt_file

# ================================== #
# PDF #
# ================================== #
def create_pdf(content):
    pdf = "hasil_modul.pdf"
    doc = SimpleDocTemplate(pdf)
    style = styles.getSampleStyleSheet()
    data = []
    for line in content.split("\n"):
        data.append(Paragraph(line, style['BodyText']))
        data.append(Spacer(1, 6))
    doc.build(data)
    return pdf

# ================================== #
# BARU: BUAT HTML PRESENTASI (Reveal.js) #
# ================================== #
def create_html(content, quiz):
    html_file = "hasil_presentasi.html"
    
    # Generate section slides
    slides_html = ""
    
    # Slide Pembuka
    slides_html += f"""
    <section data-background-color="#{THEME['primary']}" style="color: white;">
        <h1 style="color: #{THEME['secondary']}; text-shadow: 2px 2px 4px rgba(0,0,0,0.4);">⛪ Presentasi Modul Ajar</h1>
        <p style="font-size: 1.2em; opacity: 0.9;">Pendidikan Agama Katolik</p>
    </section>
    """
    
    # Slide dari AI
    slides = content.split("Slide")
    for item in slides:
        if len(item) < 10:
            continue
        parts = item.split("\n")
        title = parts[0].strip(" :0123456789")
        
        bullets = ""
        for line in parts[1:]:
            if line.strip():
                # Bersihkan pointer bawaan AI jika ada
                clean_line = line.strip().lstrip("-*• ").strip()
                bullets += f"<li style='margin-bottom: 15px;'>{clean_line}</li>"
        
        slides_html += f"""
        <section data-background-color="#f8f9fa">
            <h2 style="color: #{THEME['primary']}; border-bottom: 3px solid #{THEME['secondary']}; padding-bottom: 10px;">{title}</h2>
            <div style="text-align: left; display: inline-block; width: 85%; margin-top: 30px;">
                <ul style="color: #333; font-size: 0.9em; line-height: 1.6;">
                    {bullets}
                </ul>
            </div>
        </section>
        """
        
    # Slide Kuis
    quiz_paragraphs = ""
    for line in quiz.split("\n"):
        if line.strip():
            quiz_paragraphs += f"<p style='text-align: left; font-size: 0.65em; margin: 5px 0;'>{line.strip()}</p>"
            
    slides_html += f"""
    <section data-background-color="#f8f9fa">
        <h2 style="color: #{THEME['primary']}; border-bottom: 3px solid #{THEME['secondary']}; padding-bottom: 10px;">📋 Kuis Interaktif</h2>
        <div style="max-height: 500px; overflow-y: auto; padding: 10px;">
            {quiz_paragraphs}
        </div>
    </section>
    """

    # Template HTML utuh dengan Reveal.js CDN & Styling Liturgi/Katolik yang elegan
    full_html = f"""
    <!doctype html>
    <html lang="id">
    <head>
        <meta charset="utf-8">
        <title>Presentasi Modul Ajar Katolik</title>
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/serif.min.css">
        <style>
            .reveal h1, .reveal h2, .reveal h3 {{ font-family: 'Georgia', serif; font-weight: bold; }}
            .reveal .controls-arrow {{ color: #{THEME['secondary']} !important; }}
            .reveal .progress span {{ background: #{THEME['secondary']} !important; }}
        </style>
    </head>
    <body>
        <div class="reveal">
            <div class="slides">
                {slides_html}
            </div>
        </div>
        <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.js"></script>
        <script>
            Reveal.initialize({{
                hash: true,
                center: true,
                transition: 'slide'
            }});
        </script>
    </body>
    </html>
    """
    
    with open(html_file, "w", encoding="utf-8") as f:
        f.write(full_html)
    return html_file

# ================================== #
# UPLOAD & RUN #
# ================================== #
uploaded = st.file_uploader("Upload Modul", type=["docx", "pdf"])
if uploaded:
    ext = os.path.splitext(uploaded.name)[1]
    if st.button("🚀 Buat Presentasi Interaktif"):
        with st.spinner("Gemini AI sedang bekerja..."):
            if ext == ".docx":
                text = read_docx(uploaded)
            else:
                text = read_pdf(uploaded)
                
            slide_content = generate_slides(text)
            quiz = generate_quiz(text)
            
            # Generate file-file output
            ppt = create_ppt(slide_content, quiz)
            pdf = create_pdf(slide_content)
            html_res = create_html(slide_content, quiz) # Panggil fungsi HTML baru
            
            st.success("Selesai! File berhasil dibuat.")
            
            # Baris Tombol Download
            col1, col2, col3 = st.columns(3)
            with col1:
                with open(ppt, "rb") as f:
                    st.download_button("⬇ Download PPT", f, "presentasi.pptx")
            with col2:
                with open(pdf, "rb") as f:
                    st.download_button("⬇ Download PDF", f, "presentasi.pdf")
            with col3:
                with open(html_res, "rb") as f:
                    st.download_button("🌐 Download HTML Interaktif", f, "presentasi.html", mime="text/html")
            
            # Tampilkan Preview HTML langsung di aplikasi Streamlit
            st.subheader("🖥️ Live Preview Presentasi")
            with open(html_res, "r", encoding="utf-8") as f:
                html_content = f.read()
            st.components.v1.html(html_content, height=550, scrolling=True)
            
            st.subheader("Preview Materi Raw")
            st.write(slide_content)
            st.subheader("Preview Kuis Raw")
            st.write(quiz)
